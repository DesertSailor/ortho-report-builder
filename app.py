import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from streamlit_cropper import st_cropper
from PIL import Image
import io
import os
import pickle
from datetime import date

# --- Formatting Helpers ---
def clean_dr_name(name):
    """Ensures name is properly capitalized and prefixed with Dr. without duplication"""
    name = name.strip()
    if not name or name.lower() == "none":
        return ""
    if name.lower().startswith("dr"):
        core = name[2:].strip().lstrip(".").strip()
        return f"Dr. {core.title()}"
    return f"Dr. {name.title()}"

st.set_page_config(page_title="Ortho Morning Report Builder", layout="wide")
st.title("🏥 Orthopedic Morning Report Builder")

# --- 1. Personal Session Isolation (Sidebar) ---
st.sidebar.header("👤 Session Privacy")
user_key = st.sidebar.text_input(
    "Enter Your Initials / Shift ID", 
    value="Resident_1", 
    help="Isolates your queue. Change this to your initials so other users on this link do not overwrite your data."
)

# Sanitize the input key to create a valid, safe local filename
safe_key = "".join([c for c in user_key if c.isalnum()]).lower()
if not safe_key:
    safe_key = "default"

BACKUP_FILE = f"morning_report_{safe_key}.pkl"

def save_local_backup():
    """Saves the current session queue to isolated local file system"""
    with open(BACKUP_FILE, 'wb') as f:
        pickle.dump(st.session_state.cases, f)

# --- 2. Session State Isolation & Dynamic Recovery Engine ---
# If switching users or loading for the first time, load the corresponding private data file
if 'current_user_token' not in st.session_state or st.session_state.current_user_token != safe_key:
    st.session_state.current_user_token = safe_key
    
    if os.path.exists(BACKUP_FILE):
        try:
            with open(BACKUP_FILE, 'rb') as f:
                st.session_state.cases = pickle.load(f)
        except Exception:
            st.session_state.cases = []
    else:
        st.session_state.cases = []
        
    st.session_state.case_counter = len(st.session_state.cases)

# Visual notification confirming the active dataset is securely isolated
if st.session_state.cases:
    st.sidebar.success(f"🔒 Isolated Session: {len(st.session_state.cases)} cases loaded for key '{safe_key}'.")
else:
    st.sidebar.info(f"🔓 Isolated Workspace: Ready for key '{safe_key}'.")

# --- 3. Global Presentation Settings ---
st.header("📅 Presentation Settings")
presentation_date = st.date_input("Date of Presentation (Calendar)", date.today())
presentation_date_str = presentation_date.strftime("%A, %B %d, %Y")

st.markdown("---")

# --- 4. Duty Team Setup ---
st.header("👥 Today's Duty Team")
with st.expander("Review Consultants & Residents on Duty", expanded=True):
    c1, c2 = st.columns(2)
    with c1:
        st.subheader("Consultants")
        
        # Base roster datasets
        eopd_ward_options = ["None", "Dr. Addisu", "Dr. Tilahun", "Dr. Samuel", "Dr. Mathias", "Dr. Tesfatsion", "Dr. Kalkidan", "Dr. Ashenafi"]
        sport_options = ["None", "Dr. Mahder", "Dr. Mamo"]
        trauma_options = ["None", "Dr. Abiy", "Dr. Milkias", "Dr. Beza", "Dr. Ibrahim"]
        
        # Spot 1 Choice
        cons_1 = st.selectbox("Consultant Spot 1 (EOPD/Ward)", eopd_ward_options, index=0)
        
        # Dynamic Mutual Exclusion: Filter out Spot 1 selection from Spot 2 options list
        eopd_ward_options_filtered = [opt for opt in eopd_ward_options if opt == "None" or opt != cons_1]
        cons_2 = st.selectbox("Consultant Spot 2 (EOPD/Ward)", eopd_ward_options_filtered, index=0)
        
        cons_3 = st.selectbox("Consultant Spot 3 (Sport)", sport_options, index=0)
        cons_4 = st.selectbox("Consultant Spot 4 (Trauma)", trauma_options, index=0)
        
        st.text_input("Consultant Spot 5 (Permanent)", value="Dr. Solomon", disabled=True)
        st.text_input("Consultant Spot 6 (Permanent)", value="Dr. Dawit", disabled=True)
    with c2:
        st.subheader("Residents")
        res1 = st.text_input("Resident 1", placeholder="Name")
        res2 = st.text_input("Resident 2", placeholder="Name")
        res3 = st.text_input("Resident 3", placeholder="Name")
        res4 = st.text_input("Resident 4", placeholder="Name")
        res5 = st.text_input("Resident 5", placeholder="Name")
        res6 = st.text_input("Resident 6", placeholder="Name")

st.markdown("---")

# --- 5. Case Entry Section ---
st.header("📝 Add Patient Case")
idx = st.session_state.case_counter

cc1, cc2 = st.columns(2)
with cc1:
    p_name = st.text_input("Patient Initials / Name", key=f"name_{idx}")
    
    # Strict numeric verification
    raw_mrn = st.text_input("MRN (Numbers Only)", key=f"raw_mrn_{idx}")
    mrn_is_valid = True
    if raw_mrn:
        if not raw_mrn.isdigit():
            st.error("❌ Strict Constraint Violation: MRN must contain numbers only. Letters or spaces are prohibited.")
            mrn_is_valid = False
        p_mrn = raw_mrn
    else:
        p_mrn = ""
        
    p_age = st.text_input("Age", key=f"age_{idx}")
    p_sex = st.selectbox("Sex", ["M", "F"], key=f"sex_{idx}")

with cc2:
    p_doi = st.date_input("Date of Injury", date.today(), key=f"doi_{idx}")
    p_doi_str = p_doi.strftime("%b %d, %Y")

    moi_options = ["RTA", "Fall from Height", "FDA", "Sports Injury", "Direct Blow", "Others"]
    selected_moi = st.selectbox("Mechanism of Injury (MOI)", moi_options, key=f"moi_sel_{idx}")
    if selected_moi == "Others":
        p_moi = st.text_input("Specify Custom MOI", key=f"moi_txt_{idx}")
    else:
        p_moi = selected_moi

    p_duration = st.text_input("Duration of Injury (e.g., 3 hrs, 2 days)", key=f"dur_{idx}")
    is_operated = st.radio("Was this case operated during the shift?", ["No", "Yes"], key=f"op_{idx}")

# Interactive Additional Entry Block
st.subheader("💡 Optional Case Specifics")
p_notes = st.text_area("Additional Remarks / Notes", placeholder="Enter specific diagnostics, laboratory tracking values, neurovascular state, or custom plan parameters...", key=f"notes_{idx}")

# Processing Multiple Picture Upload Blocks
cropped_pre_list = []
cropped_post_list = []

st.subheader("📸 Radiograph Batch Processing")
img_col1, img_col2 = st.columns(2)

with img_col1:
    pre_files = st.file_uploader("Upload Injury / Pre-Op Images", type=["jpg", "jpeg", "png"], accept_multiple_files=True, key=f"pre_files_{idx}")
    if pre_files:
        for f_idx, f in enumerate(pre_files):
            st.write(f"🔍 **Pre-Op Image [{f_idx + 1}]:**")
            raw_img = Image.open(f)
            c_img = st_cropper(raw_img, realtime_update=False, box_color='#FF0000', aspect_ratio=None, key=f"crop_pre_{idx}_{f_idx}")
            cropped_pre_list.append(c_img)

with img_col2:
    if is_operated == "Yes":
        post_files = st.file_uploader("Upload Post-Op / Fixation Images", type=["jpg", "jpeg", "png"], accept_multiple_files=True, key=f"post_files_{idx}")
        if post_files:
            for f_idx, f in enumerate(post_files):
                st.write(f"🔍 **Post-Op Image [{f_idx + 1}]:**")
                raw_img = Image.open(f)
                c_img = st_cropper(raw_img, realtime_update=False, box_color='#00FF00', aspect_ratio=None, key=f"crop_post_{idx}_{f_idx}")
                cropped_post_list.append(c_img)

# Commit Entries to Session Cache & File Backup
if st.button("➕ Save This Case"):
    if not mrn_is_valid:
        st.error("Please provide a valid numeric MRN before saving.")
    elif p_name or p_mrn:
        new_case = {
            "name": p_name, "mrn": p_mrn, "age": p_age, "sex": p_sex, "moi": p_moi, "duration": p_duration,
            "doi": p_doi_str,
            "operated": is_operated,
            "notes": p_notes.strip(),
            "pre_imgs": cropped_pre_list.copy(),
            "post_imgs": cropped_post_list.copy() if is_operated == "Yes" else []
        }
        st.session_state.cases.append(new_case)
        st.session_state.case_counter += 1
        
        # Flash write directly to local file backup
        save_local_backup()
        st.rerun()
    else:
        st.error("Please provide at least a Patient Identifier or MRN before saving.")

st.markdown("---")

# --- 6. Queue Display & File Processing Engine ---
if st.session_state.cases:
    st.header(f"📋 Presentation Queue ({len(st.session_state.cases)} Cases Added for Session '{safe_key}')")
    for i, c in enumerate(st.session_state.cases):
        tag = "🔴 Operated" if c['operated'] == "Yes" else "🟢 Conservative"
        st.write(f"**Case {i+1}:** {c['name']} ({c['age']}/{c['sex']}) — MRN: {c['mrn']} | {tag}")
        
    if st.button("🗑️ Reset Entire Queue"):
        st.session_state.cases = []
        st.session_state.case_counter = 0
        if os.path.exists(BACKUP_FILE):
            os.remove(BACKUP_FILE)
        st.rerun()

    st.subheader("🚀 Export Slide Deck")
    if st.button("Compile PowerPoint Presentation"):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        # --- Slide 1: Cover Presentation ---
        slide1 = prs.slides.add_slide(blank_layout)
        tb = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = f"Duty activity of {presentation_date_str}"
        p_title.font.size = Pt(44)
        p_title.font.bold = True
        
        # Compile structured names smoothly
        raw_seniors = [cons_1, cons_2, cons_3, cons_4, "Dr. Solomon", "Dr. Dawit"]
        seniors = [clean_dr_name(s) for s in raw_seniors if s and s.lower() != "none"]
        residents = [clean_dr_name(r) for r in [res1, res2, res3, res4, res5, res6] if r.strip()]
        
        p_team = tf.add_paragraph()
        p_team.text = f"\n🔹 Consultants on Duty:\n{', '.join(seniors) if seniors else 'None Specified'}\n\n" \
                     f"🔹 Residents on Duty:\n{', '.join(residents) if residents else 'None Specified'}"
        p_team.font.size = Pt(22)

        # --- Priority Queue Ordering Logic (Operated Cases First) ---
        priority_sorted_cases = [c for c in st.session_state.cases if c['operated'] == "Yes"] + \
                                [c for c in st.session_state.cases if c['operated'] == "No"]

        # --- Case Document Generation Loop ---
        for case_idx, c in enumerate(priority_sorted_cases, start=1):
            all_imgs = []
            for img in c['pre_imgs']:
                all_imgs.append(img)
            for img in c['post_imgs']:
                all_imgs.append(img)
                
            # Paginate images to a maximum of 2 side-by-side per slide
            img_chunks = [all_imgs[i:i + 2] for i in range(0, len(all_imgs), 2)]
            if not img_chunks:
                img_chunks = [[]]
                
            for chunk_idx, chunk in enumerate(img_chunks):
                slide_case = prs.slides.add_slide(blank_layout)
                
                # Render text descriptors strictly on the first slide for the case
                if chunk_idx == 0:
                    # TITLE BOX: Demographic layout line
                    title_box = slide_case.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
                    tf_title = title_box.text_frame
                    p_t = tf_title.paragraphs[0]
                    p_t.text = f"{case_idx}. {c['name']}  |  {c['age']}/{c['sex']}  |  MRN: {c['mrn']}"
                    p_t.font.size = Pt(32)
                    p_t.font.bold = True
                    
                    # DESCRIPTION & REMARKS BOX
                    desc_box = slide_case.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.9))
                    tf_desc = desc_box.text_frame
                    tf_desc.word_wrap = True
                    
                    p_d = tf_desc.paragraphs[0]
                    p_d.text = f"MOI: {c['moi']}   |   DOI: {c['doi']}   |   Duration: {c['duration']}"
                    p_d.font.size = Pt(18)
                    p_d.font.color.rgb = RGBColor(60, 60, 60)
                    
                    # Custom user instructions string printed cleanly without label prefixes
                    if c['notes']:
                        p_n = tf_desc.add_paragraph()
                        p_n.text = f"{c['notes']}"
                        p_n.font.size = Pt(15)
                        p_n.font.italic = True
                        p_n.font.color.rgb = RGBColor(90, 90, 90)
                        
                    img_top = Inches(2.3)
                else:
                    # Clear multi-slide option layout: Images ONLY for sequential slides
                    img_top = Inches(1.0)
                
                # Rendering logic based on chunk allocations
                if chunk:
                    if len(chunk) == 1:
                        img_obj = chunk[0]
                        img_buf = io.BytesIO()
                        img_obj.save(img_buf, format="PNG")
                        img_buf.seek(0)
                        slide_case.shapes.add_picture(img_buf, Inches(3.766), img_top, width=Inches(5.8))
                    elif len(chunk) == 2:
                        img_buf1 = io.BytesIO()
                        chunk[0].save(img_buf1, format="PNG")
                        img_buf1.seek(0)
                        slide_case.shapes.add_picture(img_buf1, Inches(0.6), img_top, width=Inches(5.8))
                        
                        img_buf2 = io.BytesIO()
                        chunk[1].save(img_buf2, format="PNG")
                        img_buf2.seek(0)
                        slide_case.shapes.add_picture(img_buf2, Inches(6.933), img_top, width=Inches(5.8))

        # Output Final Presentation Binary
        final_ppt_buf = io.BytesIO()
        prs.save(final_ppt_buf)
        final_ppt_buf.seek(0)
        
        st.download_button(
            label="💾 Save PowerPoint File to Device",
            data=final_ppt_buf,
            file_name=f"Morning_Report_{presentation_date.strftime('%Y_%m_%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
                st.session_state.cases = pickle.load(f)
        except Exception:
            st.session_state.cases = []
    else:
        st.session_state.cases = []

if 'case_counter' not in st.session_state:
    st.session_state.case_counter = len(st.session_state.cases)

st.set_page_config(page_title="Ortho Morning Report Builder", layout="wide")
st.title("🏥 Orthopedic Morning Report Builder")

# Visual feedback tracker for active database sessions
if st.session_state.cases:
    st.success(f"💾 Local persistent backup active: {len(st.session_state.cases)} cases safely cached on disk.")

# --- 1. Global Presentation Settings ---
st.header("📅 Presentation Settings")
presentation_date = st.date_input("Date of Presentation (Calendar)", date.today())
presentation_date_str = presentation_date.strftime("%A, %B %d, %Y")

st.markdown("---")

# --- 2. Duty Team Setup ---
st.header("👥 Today's Duty Team")
with st.expander("Review Consultants & Residents on Duty", expanded=True):
    c1, c2 = st.columns(2)
    with c1:
        st.subheader("Consultants")
        
        # Base roster datasets
        eopd_ward_options = ["None", "Dr. Addisu", "Dr. Tilahun", "Dr. Samuel", "Dr. Mathias", "Dr. Tesfatsion", "Dr. Kalkidan", "Dr. Ashenafi"]
        sport_options = ["None", "Dr. Mahder", "Dr. Mamo"]
        trauma_options = ["None", "Dr. Abiy", "Dr. Milkias", "Dr. Beza", "Dr. Ibrahim"]
        
        # Spot 1 Choice
        cons_1 = st.selectbox("Consultant Spot 1 (EOPD/Ward)", eopd_ward_options, index=0)
        
        # Dynamic Mutual Exclusion: Filter out Spot 1 selection from Spot 2 options list
        eopd_ward_options_filtered = [opt for opt in eopd_ward_options if opt == "None" or opt != cons_1]
        cons_2 = st.selectbox("Consultant Spot 2 (EOPD/Ward)", eopd_ward_options_filtered, index=0)
        
        cons_3 = st.selectbox("Consultant Spot 3 (Sport)", sport_options, index=0)
        cons_4 = st.selectbox("Consultant Spot 4 (Trauma)", trauma_options, index=0)
        
        st.text_input("Consultant Spot 5 (Permanent)", value="Dr. Solomon", disabled=True)
        st.text_input("Consultant Spot 6 (Permanent)", value="Dr. Dawit", disabled=True)
    with c2:
        st.subheader("Residents")
        res1 = st.text_input("Resident 1", placeholder="Name")
        res2 = st.text_input("Resident 2", placeholder="Name")
        res3 = st.text_input("Resident 3", placeholder="Name")
        res4 = st.text_input("Resident 4", placeholder="Name")
        res5 = st.text_input("Resident 5", placeholder="Name")
        res6 = st.text_input("Resident 6", placeholder="Name")

st.markdown("---")

# --- 3. Case Entry Section ---
st.header("📝 Add Patient Case")
idx = st.session_state.case_counter

cc1, cc2 = st.columns(2)
with cc1:
    p_name = st.text_input("Patient Initials / Name", key=f"name_{idx}")
    
    # Strict numeric verification
    raw_mrn = st.text_input("MRN (Numbers Only)", key=f"raw_mrn_{idx}")
    mrn_is_valid = True
    if raw_mrn:
        if not raw_mrn.isdigit():
            st.error("❌ Strict Constraint Violation: MRN must contain numbers only. Letters or spaces are prohibited.")
            mrn_is_valid = False
        p_mrn = raw_mrn
    else:
        p_mrn = ""
        
    p_age = st.text_input("Age", key=f"age_{idx}")
    p_sex = st.selectbox("Sex", ["M", "F"], key=f"sex_{idx}")

with cc2:
    p_doi = st.date_input("Date of Injury", date.today(), key=f"doi_{idx}")
    p_doi_str = p_doi.strftime("%b %d, %Y")

    moi_options = ["RTA", "Fall from Height", "FDA", "Sports Injury", "Direct Blow", "Others"]
    selected_moi = st.selectbox("Mechanism of Injury (MOI)", moi_options, key=f"moi_sel_{idx}")
    if selected_moi == "Others":
        p_moi = st.text_input("Specify Custom MOI", key=f"moi_txt_{idx}")
    else:
        p_moi = selected_moi

    p_duration = st.text_input("Duration of Injury (e.g., 3 hrs, 2 days)", key=f"dur_{idx}")
    is_operated = st.radio("Was this case operated during the shift?", ["No", "Yes"], key=f"op_{idx}")

# Interactive Additional Entry Block
st.subheader("💡 Optional Case Specifics")
p_notes = st.text_area("Additional Remarks / Notes", placeholder="Enter specific diagnostics, laboratory tracking values, neurovascular state, or custom plan parameters...", key=f"notes_{idx}")

# Processing Multiple Picture Upload Blocks
cropped_pre_list = []
cropped_post_list = []

st.subheader("📸 Radiograph Batch Processing")
img_col1, img_col2 = st.columns(2)

with img_col1:
    pre_files = st.file_uploader("Upload Injury / Pre-Op Images", type=["jpg", "jpeg", "png"], accept_multiple_files=True, key=f"pre_files_{idx}")
    if pre_files:
        for f_idx, f in enumerate(pre_files):
            st.write(f"🔍 **Pre-Op Image [{f_idx + 1}]:**")
            raw_img = Image.open(f)
            c_img = st_cropper(raw_img, realtime_update=False, box_color='#FF0000', aspect_ratio=None, key=f"crop_pre_{idx}_{f_idx}")
            cropped_pre_list.append(c_img)

with img_col2:
    if is_operated == "Yes":
        post_files = st.file_uploader("Upload Post-Op / Fixation Images", type=["jpg", "jpeg", "png"], accept_multiple_files=True, key=f"post_files_{idx}")
        if post_files:
            for f_idx, f in enumerate(post_files):
                st.write(f"🔍 **Post-Op Image [{f_idx + 1}]:**")
                raw_img = Image.open(f)
                c_img = st_cropper(raw_img, realtime_update=False, box_color='#00FF00', aspect_ratio=None, key=f"crop_post_{idx}_{f_idx}")
                cropped_post_list.append(c_img)

# Commit Entries to Session Cache & File Backup
if st.button("➕ Save This Case"):
    if not mrn_is_valid:
        st.error("Please provide a valid numeric MRN before saving.")
    elif p_name or p_mrn:
        new_case = {
            "name": p_name, "mrn": p_mrn, "age": p_age, "sex": p_sex, "moi": p_moi, "duration": p_duration,
            "doi": p_doi_str,
            "operated": is_operated,
            "notes": p_notes.strip(),
            "pre_imgs": cropped_pre_list.copy(),
            "post_imgs": cropped_post_list.copy() if is_operated == "Yes" else []
        }
        st.session_state.cases.append(new_case)
        st.session_state.case_counter += 1
        
        # Flash write directly to local file backup
        save_local_backup()
        st.invalidate_pages() if hasattr(st, "invalidate_pages") else None
        st.rerun()
    else:
        st.error("Please provide at least a Patient Identifier or MRN before saving.")

st.markdown("---")

# --- 4. Queue Display & File Processing Engine ---
if st.session_state.cases:
    st.header(f"📋 Presentation Queue ({len(st.session_state.cases)} Cases Added)")
    for i, c in enumerate(st.session_state.cases):
        tag = "🔴 Operated" if c['operated'] == "Yes" else "🟢 Conservative"
        st.write(f"**Case {i+1}:** {c['name']} ({c['age']}/{c['sex']}) — MRN: {c['mrn']} | {tag}")
        
    if st.button("🗑️ Reset Entire Queue"):
        st.session_state.cases = []
        st.session_state.case_counter = 0
        if os.path.exists(BACKUP_FILE):
            os.remove(BACKUP_FILE)
        st.rerun()

    st.subheader("🚀 Export Slide Deck")
    if st.button("Compile PowerPoint Presentation"):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        # --- Slide 1: Cover Presentation ---
        slide1 = prs.slides.add_slide(blank_layout)
        tb = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = f"Duty activity of {presentation_date_str}"
        p_title.font.size = Pt(44)
        p_title.font.bold = True
        
        # Compile structured names smoothly
        raw_seniors = [cons_1, cons_2, cons_3, cons_4, "Dr. Solomon", "Dr. Dawit"]
        seniors = [clean_dr_name(s) for s in raw_seniors if s and s.lower() != "none"]
        residents = [clean_dr_name(r) for r in [res1, res2, res3, res4, res5, res6] if r.strip()]
        
        p_team = tf.add_paragraph()
        p_team.text = f"\n🔹 Consultants on Duty:\n{', '.join(seniors) if seniors else 'None Specified'}\n\n" \
                     f"🔹 Residents on Duty:\n{', '.join(residents) if residents else 'None Specified'}"
        p_team.font.size = Pt(22)

        # --- Priority Queue Ordering Logic (Operated Cases First) ---
        priority_sorted_cases = [c for c in st.session_state.cases if c['operated'] == "Yes"] + \
                                [c for c in st.session_state.cases if c['operated'] == "No"]

        # --- Case Document Generation Loop ---
        for case_idx, c in enumerate(priority_sorted_cases, start=1):
            all_imgs = []
            for img in c['pre_imgs']:
                all_imgs.append(img)
            for img in c['post_imgs']:
                all_imgs.append(img)
                
            # Paginate images to a maximum of 2 side-by-side per slide
            img_chunks = [all_imgs[i:i + 2] for i in range(0, len(all_imgs), 2)]
            if not img_chunks:
                img_chunks = [[]]
                
            for chunk_idx, chunk in enumerate(img_chunks):
                slide_case = prs.slides.add_slide(blank_layout)
                
                # Render text descriptors strictly on the first slide for the case
                if chunk_idx == 0:
                    # TITLE BOX: Demographic layout line
                    title_box = slide_case.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
                    tf_title = title_box.text_frame
                    p_t = tf_title.paragraphs[0]
                    p_t.text = f"{case_idx}. {c['name']}  |  {c['age']}/{c['sex']}  |  MRN: {c['mrn']}"
                    p_t.font.size = Pt(32)
                    p_t.font.bold = True
                    
                    # DESCRIPTION & REMARKS BOX
                    desc_box = slide_case.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.9))
                    tf_desc = desc_box.text_frame
                    tf_desc.word_wrap = True
                    
                    p_d = tf_desc.paragraphs[0]
                    p_d.text = f"MOI: {c['moi']}   |   DOI: {c['doi']}   |   Duration: {c['duration']}"
                    p_d.font.size = Pt(18)
                    p_d.font.color.rgb = RGBColor(60, 60, 60)
                    
                    # Custom user instructions string printed cleanly without label prefixes
                    if c['notes']:
                        p_n = tf_desc.add_paragraph()
                        p_n.text = f"{c['notes']}"
                        p_n.font.size = Pt(15)
                        p_n.font.italic = True
                        p_n.font.color.rgb = RGBColor(90, 90, 90)
                        
                    img_top = Inches(2.3)
                else:
                    # Clear multi-slide option layout: Images ONLY for sequential slides
                    img_top = Inches(1.0)
                
                # Rendering logic based on chunk allocations
                if chunk:
                    if len(chunk) == 1:
                        img_obj = chunk[0]
                        img_buf = io.BytesIO()
                        img_obj.save(img_buf, format="PNG")
                        img_buf.seek(0)
                        slide_case.shapes.add_picture(img_buf, Inches(3.766), img_top, width=Inches(5.8))
                    elif len(chunk) == 2:
                        img_buf1 = io.BytesIO()
                        chunk[0].save(img_buf1, format="PNG")
                        img_buf1.seek(0)
                        slide_case.shapes.add_picture(img_buf1, Inches(0.6), img_top, width=Inches(5.8))
                        
                        img_buf2 = io.BytesIO()
                        chunk[1].save(img_buf2, format="PNG")
                        img_buf2.seek(0)
                        slide_case.shapes.add_picture(img_buf2, Inches(6.933), img_top, width=Inches(5.8))

        # Output Final Presentation Binary
        final_ppt_buf = io.BytesIO()
        prs.save(final_ppt_buf)
        final_ppt_buf.seek(0)
        
        st.download_button(
            label="💾 Save PowerPoint File to Device",
            data=final_ppt_buf,
            file_name=f"Morning_Report_{presentation_date.strftime('%Y_%m_%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
